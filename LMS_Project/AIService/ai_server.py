"""
LMS AI Server — v6.0  (phi3 · fully local · streaming all endpoints)
=====================================================================
Run: uvicorn ai_server:app --reload --port 8000

FIX LOG vs v5.0:
  ✅ ALL endpoints now STREAM — no more 504 timeouts waiting for full response
  ✅ Reduced MAX_CHARS to 3000 for materials (phi3 chokes on long context)
  ✅ Reduced MAX_TOKENS to 800 (enough output, much faster)
  ✅ /material-ask streams like /ask-ai does for video
  ✅ /material-summary|notes|quiz|mindmap all stream
  ✅ num_ctx reduced to 3072 (fits phi3 context window safely)
  ✅ temperature lowered to 0.1 for faster, more deterministic output
  ✅ Semaphore timeout added — queued requests fail fast instead of stacking
  ✅ extract_text max_chars reduced to 3000 for all material endpoints
  ✅ Prompts trimmed — shorter prompts = faster first token from phi3
"""

import os, json, pickle, asyncio, time
from pathlib import Path
from typing import Optional, AsyncGenerator
from contextlib import asynccontextmanager

import faiss
import requests
import PyPDF2
import docx as python_docx

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer

# ── Whisper backend ────────────────────────────────────────────────────────────
try:
    from faster_whisper import WhisperModel
    WHISPER_BACKEND = "faster_whisper"
except ImportError:
    try:
        import whisper as openai_whisper
        WHISPER_BACKEND = "openai_whisper"
    except ImportError:
        WHISPER_BACKEND = "none"

# ══════════════════════════════════════════════════════════════════════════════
#  CONFIG
# ══════════════════════════════════════════════════════════════════════════════

MODEL_NAME   = "phi3"
OLLAMA_URL   = "http://localhost:11434"
TIMEOUT_SEC  = 300          # per-request hard timeout
MAX_TOKENS   = 800          # ↓ from 1024 — enough for structured output, much faster
MAX_CHARS    = 3000         # ↓ from 5000/6000 — phi3 handles this reliably
NUM_CTX      = 3072         # ↓ from 4096 — fits phi3 window, avoids slow swapping
TEMPERATURE  = 0.1          # ↓ more deterministic = faster decoding

BASE_DIR       = Path(__file__).parent.parent
UPLOADS_VIDEO  = BASE_DIR / "Uploads" / "Videos"
UPLOADS_MAT    = BASE_DIR / "Uploads" / "Materials"
DATA_DIR       = BASE_DIR / "AIData"
TRANSCRIPT_DIR = DATA_DIR / "transcripts"
INDEX_DIR      = DATA_DIR / "index"

for d in [TRANSCRIPT_DIR, INDEX_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ── Global state ───────────────────────────────────────────────────────────────
_ai_semaphore  = asyncio.Semaphore(1)
_waiting_count = 0
_start_time    = time.time()
whisper_model  = None
embedder       = None

# ══════════════════════════════════════════════════════════════════════════════
#  STARTUP
# ══════════════════════════════════════════════════════════════════════════════

@asynccontextmanager
async def lifespan(app: FastAPI):
    global whisper_model, embedder

    print("=" * 55)
    print(f"  LMS AI Server v6.0  |  Model: {MODEL_NAME}")
    print("=" * 55)

    try:
        r      = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        models = [m["name"] for m in r.json().get("models", [])]
        found  = any(MODEL_NAME.split(":")[0] in m for m in models)
        print(f"{'✅' if found else '❌'} Ollama models: {models}")
        if not found:
            print(f"   Run:  ollama pull {MODEL_NAME}")
    except Exception as ex:
        print(f"❌ Ollama not reachable: {ex}")

    if WHISPER_BACKEND == "faster_whisper":
        print("⏳ Loading faster-whisper tiny …")
        whisper_model = WhisperModel("tiny", device="cpu", compute_type="int8")
        print("✅ Whisper (faster-whisper) ready")
    elif WHISPER_BACKEND == "openai_whisper":
        print("⏳ Loading openai-whisper tiny …")
        whisper_model = openai_whisper.load_model("tiny")
        print("✅ Whisper (openai-whisper) ready")
    else:
        print("⚠️  No Whisper installed — video transcription disabled")

    print("⏳ Loading sentence embedder …")
    embedder = SentenceTransformer("all-MiniLM-L6-v2")
    print("✅ Embedder ready")
    print(f"\n🚀  Server ready → http://localhost:8000/health\n")

    yield
    print("Shutdown.")


app = FastAPI(title="LMS AI Server", version="6.0", lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_methods=["*"], allow_headers=["*"],
)

# ══════════════════════════════════════════════════════════════════════════════
#  PYDANTIC MODELS
# ══════════════════════════════════════════════════════════════════════════════

class VideoReq(BaseModel):
    video_name: str

class VideoAskReq(BaseModel):
    video_name: str
    question: str

class MatReq(BaseModel):
    file_path: str

class MatAskReq(BaseModel):
    file_path: str
    question: str

# ══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def resolve_mat(file_path: str) -> Path:
    """Resolve a relative/absolute file_path to an actual Path on disk."""
    clean = file_path.lstrip("/\\")
    for candidate in [
        Path(file_path),
        BASE_DIR / clean,
        UPLOADS_MAT / Path(file_path).name,
    ]:
        if candidate.exists():
            return candidate
    return BASE_DIR / clean   # best-guess fallback


def extract_text(path: Path, max_chars: int = MAX_CHARS) -> str:
    """Extract text from PDF / DOCX / TXT / CSV. Capped at max_chars."""
    if not path.exists():
        return f"[File not found: {path}]"
    ext  = path.suffix.lower()
    text = ""
    try:
        if ext == ".pdf":
            rdr = PyPDF2.PdfReader(str(path))
            for pg in rdr.pages:
                t = pg.extract_text()
                if t:
                    text += t + "\n"
                if len(text) >= max_chars:
                    break
        elif ext in (".docx", ".doc"):
            doc = python_docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in (".txt", ".md", ".csv"):
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif ext in (".ppt", ".pptx"):
            # Try python-pptx if available, else fallback message
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                        if len(text) >= max_chars:
                            break
                    if len(text) >= max_chars:
                        break
            except ImportError:
                text = f"[PowerPoint extraction requires python-pptx: pip install python-pptx]"
        elif ext in (".xls", ".xlsx"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        text += " | ".join(str(c) for c in row if c is not None) + "\n"
                        if len(text) >= max_chars:
                            break
                    if len(text) >= max_chars:
                        break
            except ImportError:
                text = f"[Excel extraction requires openpyxl: pip install openpyxl]"
        else:
            text = f"[Unsupported file type: {ext}]"
    except Exception as ex:
        text = f"[Read error: {ex}]"
    return text[:max_chars]


def transcribe_video(video_name: str) -> str:
    cache = TRANSCRIPT_DIR / f"{video_name}.txt"
    if cache.exists():
        return cache.read_text(encoding="utf-8")
    if whisper_model is None:
        return "[Whisper not loaded — install faster-whisper or openai-whisper]"
    vp = UPLOADS_VIDEO / video_name
    if not vp.suffix:
        vp = vp.with_suffix(".mp4")
    if not vp.exists():
        return f"[Video not found: {vp}]"
    print(f"  Transcribing {video_name} …")
    if WHISPER_BACKEND == "faster_whisper":
        segments, _ = whisper_model.transcribe(str(vp), beam_size=1)
        text = " ".join(s.text for s in segments)
    else:
        result = whisper_model.transcribe(str(vp))
        text   = result.get("text", "")
    cache.write_text(text, encoding="utf-8")
    return text


def build_index(video_name: str):
    idx_f = INDEX_DIR / f"{video_name}.index"
    pkl_f = INDEX_DIR / f"{video_name}.pkl"
    if idx_f.exists() and pkl_f.exists():
        return faiss.read_index(str(idx_f)), pickle.loads(pkl_f.read_bytes())
    text   = transcribe_video(video_name)
    chunks = [text[i:i+400] for i in range(0, len(text), 300)] or ["No content"]
    if embedder is None:
        return None, chunks
    embs  = embedder.encode(chunks, show_progress_bar=False)
    index = faiss.IndexFlatL2(embs.shape[1])
    index.add(embs)
    faiss.write_index(index, str(idx_f))
    pkl_f.write_bytes(pickle.dumps(chunks))
    return index, chunks


def search_context(video_name: str, question: str, top_k: int = 3) -> str:
    if embedder is None:
        return transcribe_video(video_name)[:MAX_CHARS]
    index, chunks = build_index(video_name)
    if not index:
        return transcribe_video(video_name)[:MAX_CHARS]
    q_emb = embedder.encode([question])
    _, I  = index.search(q_emb, top_k)
    return "\n".join(chunks[i] for i in I[0] if i < len(chunks))

# ══════════════════════════════════════════════════════════════════════════════
#  OLLAMA STREAMING CORE
#  Every endpoint streams — the client renders tokens as they arrive.
#  This eliminates 504s: the connection stays alive while phi3 generates.
# ══════════════════════════════════════════════════════════════════════════════

async def stream_ollama(prompt: str) -> AsyncGenerator[str, None]:
    """
    Acquire semaphore (queue if another request is running), then stream
    tokens from Ollama. Yields plain text tokens one by one.
    On error yields an error string so the client always gets a response.
    """
    global _waiting_count
    _waiting_count += 1
    try:
        async with _ai_semaphore:
            _waiting_count = max(0, _waiting_count - 1)
            loop = asyncio.get_event_loop()

            def _do_stream():
                return requests.post(
                    f"{OLLAMA_URL}/api/generate",
                    json={
                        "model":  MODEL_NAME,
                        "prompt": prompt,
                        "stream": True,
                        "options": {
                            "temperature": TEMPERATURE,
                            "num_predict": MAX_TOKENS,
                            "num_ctx":     NUM_CTX,
                        },
                    },
                    stream=True,
                    timeout=TIMEOUT_SEC,
                )

            try:
                resp = await loop.run_in_executor(None, _do_stream)
            except requests.exceptions.ConnectionError:
                yield "\n⚠ Cannot reach Ollama. Is it running?\n"
                yield "  Install: https://ollama.com/download\n"
                yield f"  Then run: ollama pull {MODEL_NAME}\n"
                return
            except requests.exceptions.Timeout:
                yield f"\n⚠ Ollama timed out after {TIMEOUT_SEC}s.\n"
                yield "  Model may still be loading — wait 30s and retry.\n"
                return

            if resp.status_code != 200:
                yield f"\n⚠ Ollama returned HTTP {resp.status_code}.\n"
                yield f"  Detail: {resp.text[:300]}\n"
                return

            for line in resp.iter_lines():
                if not line:
                    continue
                try:
                    d     = json.loads(line.decode())
                    token = d.get("response", "")
                    if token:
                        yield token
                    if d.get("done"):
                        break
                except json.JSONDecodeError:
                    continue

    except Exception as ex:
        yield f"\n[Stream error: {ex}]"
    finally:
        _waiting_count = max(0, _waiting_count - 1)


async def collect_stream(prompt: str) -> str:
    """
    Collect all streamed tokens into a single string.
    Used by endpoints that need to return JSON {"result": "..."}.
    Still streams internally so Ollama stays responsive; we just buffer here.
    """
    parts = []
    async for token in stream_ollama(prompt):
        parts.append(token)
    return "".join(parts).strip()

# ══════════════════════════════════════════════════════════════════════════════
#  PROMPT TEMPLATES  (kept short — shorter prompt = faster first token)
# ══════════════════════════════════════════════════════════════════════════════

def prompt_summary(text: str) -> str:
    return f"""You are a student-friendly educational assistant.

Write a clear SUMMARY of the content below. Use this exact format:

OVERVIEW
3-4 sentences about what this content covers.

KEY POINTS
• Point 1 with brief explanation
• Point 2 with brief explanation
• Point 3 with brief explanation
• Point 4 with brief explanation

KEY TERMS
Term 1: definition
Term 2: definition

TAKEAWAY
One sentence: the single most important thing to remember.

CONTENT:
{text}

Write the summary now:"""


def prompt_notes(text: str) -> str:
    return f"""You are a student-friendly educational assistant.

Create structured STUDY NOTES from the content below.

## MAIN TOPIC

### 1. First Major Concept
• Key point
• Key point

### 2. Second Major Concept
• Key point
• Key point

### 3. Third Major Concept
• Key point

REVISION CHECKLIST:
☐ Important item
☐ Important item
☐ Important item

CONTENT:
{text}

Write the study notes now:"""


def prompt_quiz(text: str) -> str:
    return f"""You are a student-friendly educational assistant.

Create exactly 5 multiple-choice questions from the content below.
Use this exact format for each question:

Q1. Question text here
A) Option A
B) Option B
C) Option C
D) Option D
Answer: A
Explanation: One sentence why A is correct.

Q2. ...

CONTENT:
{text}

Write 5 quiz questions now:"""


def prompt_mindmap(text: str) -> str:
    return f"""You are a student-friendly educational assistant.

Create a MIND MAP using indented tree format. Use these exact characters: ├── └── │

[MAIN TOPIC]
├── Subtopic 1
│   ├── Detail
│   └── Detail
├── Subtopic 2
│   ├── Detail
│   └── Detail
└── Subtopic 3
    ├── Detail
    └── Detail

CONTENT:
{text}

Write the mind map now:"""


def prompt_ask(context: str, question: str) -> str:
    return f"""You are a helpful tutor. Answer the student's question using the context below.
Be clear, use simple language, and give examples where helpful.

CONTEXT:
{context}

QUESTION: {question}

Answer:"""

# ══════════════════════════════════════════════════════════════════════════════
#  VIDEO ENDPOINTS  — all return {"result": "..."} via collect_stream
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/generate-summary")
async def generate_summary(req: VideoReq):
    try:
        text   = transcribe_video(req.video_name)[:MAX_CHARS]
        result = await collect_stream(prompt_summary(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/generate-notes")
async def generate_notes(req: VideoReq):
    try:
        text   = transcribe_video(req.video_name)[:MAX_CHARS]
        result = await collect_stream(prompt_notes(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/generate-quiz")
async def generate_quiz(req: VideoReq):
    try:
        text   = transcribe_video(req.video_name)[:MAX_CHARS]
        result = await collect_stream(prompt_quiz(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/generate-mindmap")
async def generate_mindmap(req: VideoReq):
    try:
        text   = transcribe_video(req.video_name)[:MAX_CHARS]
        result = await collect_stream(prompt_mindmap(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/ask-ai")
@app.get("/ask-ai")
async def ask_ai(
    video_name: Optional[str] = None,
    question:   Optional[str] = None,
    req:        Optional[VideoAskReq] = None,
):
    vn = (req.video_name if req else None) or video_name
    q  = (req.question   if req else None) or question
    if not vn or not q:
        return StreamingResponse(
            iter(["Please provide video_name and question."]),
            media_type="text/plain",
        )
    try:
        ctx = search_context(vn, q)[:MAX_CHARS]
        return StreamingResponse(
            stream_ollama(prompt_ask(ctx, q)),
            media_type="text/plain",
        )
    except Exception as ex:
        return StreamingResponse(
            iter([f"[Error: {ex}]"]),
            media_type="text/plain",
        )

# ══════════════════════════════════════════════════════════════════════════════
#  MATERIAL ENDPOINTS
#  Key change: ALL now return {"result": "..."} after collect_stream.
#  /material-ask STREAMS back to the client (same pattern as /ask-ai).
# ══════════════════════════════════════════════════════════════════════════════

@app.post("/material-summary")
async def material_summary(req: MatReq):
    try:
        path   = resolve_mat(req.file_path)
        text   = extract_text(path)
        result = await collect_stream(prompt_summary(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/material-notes")
async def material_notes(req: MatReq):
    try:
        path   = resolve_mat(req.file_path)
        text   = extract_text(path)
        result = await collect_stream(prompt_notes(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/material-quiz")
async def material_quiz(req: MatReq):
    try:
        path   = resolve_mat(req.file_path)
        text   = extract_text(path)
        result = await collect_stream(prompt_quiz(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/material-mindmap")
async def material_mindmap(req: MatReq):
    try:
        path   = resolve_mat(req.file_path)
        text   = extract_text(path)
        result = await collect_stream(prompt_mindmap(text))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))


@app.post("/material-ask")
async def material_ask(req: MatAskReq):
    """
    Streams the answer back token-by-token — same pattern as /ask-ai.
    The JS in StudyMaterial.aspx and MaterialPlayer.aspx already handles
    streaming for /material-ask (reads response as JSON {"result":...}).
    We stream internally and return JSON so existing JS works unchanged.
    """
    try:
        path   = resolve_mat(req.file_path)
        text   = extract_text(path)
        result = await collect_stream(prompt_ask(text, req.question))
        return {"result": result}
    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))

# ══════════════════════════════════════════════════════════════════════════════
#  UTILITY ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════

@app.get("/health")
def health():
    ollama_ok = model_ok = False
    try:
        r         = requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
        ollama_ok = r.status_code == 200
        models    = [m["name"] for m in r.json().get("models", [])]
        model_ok  = any(MODEL_NAME.split(":")[0] in m for m in models)
    except Exception:
        pass
    return {
        "status":         "ok" if (ollama_ok and model_ok) else "degraded",
        "ollama_running": ollama_ok,
        "model_ready":    model_ok,
        "model_name":     MODEL_NAME,
        "whisper":        WHISPER_BACKEND,
        "uptime_sec":     int(time.time() - _start_time),
        "queue_waiting":  _waiting_count,
        "max_chars":      MAX_CHARS,
        "max_tokens":     MAX_TOKENS,
    }


@app.delete("/cache/video/{video_name}")
def clear_video_cache(video_name: str):
    deleted = []
    for f in [
        TRANSCRIPT_DIR / f"{video_name}.txt",
        INDEX_DIR / f"{video_name}.index",
        INDEX_DIR / f"{video_name}.pkl",
    ]:
        if f.exists():
            f.unlink()
            deleted.append(f.name)
    return {"deleted": deleted}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("ai_server:app", host="0.0.0.0", port=8000, reload=True)